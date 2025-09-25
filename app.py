from flask import Flask, request, send_file, jsonify
from werkzeug.utils import secure_filename
import os
import tempfile
from pdf2image import convert_from_path
from pdf2docx import Converter
import pandas as pd
from pptx import Presentation
import zipfile
from PIL import Image
import io

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB max file size

UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'output'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

ALLOWED_EXTENSIONS = {'pdf', 'jpg', 'jpeg', 'png', 'doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def cleanup_files(files):
    for file_path in files:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception as e:
            print(f"Error cleaning up file {file_path}: {e}")

@app.route('/api/convert', methods=['POST'])
def convert_file():
    temp_files = []
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400

        file = request.files['file']
        conversion_type = request.form.get('conversion_type', 'pdf-to-jpg')

        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400

        if not allowed_file(file.filename):
            return jsonify({'error': 'File type not allowed'}), 400

        filename = secure_filename(file.filename)
        input_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(input_path)
        temp_files.append(input_path)

        if conversion_type == 'pdf-to-jpg':
            output_path = convert_pdf_to_jpg(input_path)
        elif conversion_type == 'pdf-to-word':
            output_path = convert_pdf_to_word(input_path)
        elif conversion_type == 'pdf-to-ppt':
            output_path = convert_pdf_to_ppt(input_path)
        elif conversion_type == 'pdf-to-excel':
            output_path = convert_pdf_to_excel(input_path)
        elif conversion_type == 'jpg-to-pdf':
            output_path = convert_image_to_pdf(input_path)
        elif conversion_type == 'word-to-pdf':
            output_path = convert_doc_to_pdf(input_path)
        elif conversion_type == 'powerpoint-to-pdf':
            output_path = convert_ppt_to_pdf(input_path)
        elif conversion_type == 'excel-to-pdf':
            output_path = convert_excel_to_pdf(input_path)
        else:
            return jsonify({'error': 'Invalid conversion type'}), 400

        temp_files.append(output_path)
        return send_file(output_path, as_attachment=True, download_name=os.path.basename(output_path))

    except Exception as e:
        cleanup_files(temp_files)
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500

def convert_pdf_to_jpg(pdf_path):
    images = convert_from_path(pdf_path, dpi=200)
    if len(images) == 1:
        output_path = os.path.join(OUTPUT_FOLDER, 'converted.jpg')
        images[0].save(output_path, 'JPEG', quality=95)
    else:
        zip_path = os.path.join(OUTPUT_FOLDER, 'converted_images.zip')
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for i, image in enumerate(images):
                img_bytes = io.BytesIO()
                image.save(img_bytes, 'JPEG', quality=95)
                zipf.writestr(f'page_{i+1}.jpg', img_bytes.getvalue())
        output_path = zip_path
    return output_path

def convert_pdf_to_word(pdf_path):
    output_path = os.path.join(OUTPUT_FOLDER, 'converted.docx')
    cv = Converter(pdf_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()
    return output_path

def convert_pdf_to_ppt(pdf_path):
    images = convert_from_path(pdf_path, dpi=150)
    output_path = os.path.join(OUTPUT_FOLDER, 'converted.pptx')
    prs = Presentation()
    for image in images:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        img_bytes = io.BytesIO()
        image.save(img_bytes, 'PNG')
        img_bytes.seek(0)
        slide.shapes.add_picture(img_bytes, 0, 0, prs.slide_width, prs.slide_height)
    prs.save(output_path)
    return output_path

def convert_pdf_to_excel(pdf_path):
    import tabula
    tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
    if not tables:
        raise Exception("No tables found in PDF")
    output_path = os.path.join(OUTPUT_FOLDER, 'converted.xlsx')
    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
        for i, table in enumerate(tables):
            sheet_name = f'Table_{i+1}'
            table.to_excel(writer, sheet_name=sheet_name, index=False)
    return output_path

def convert_image_to_pdf(image_path):
    output_path = os.path.join(OUTPUT_FOLDER, 'converted.pdf')
    image = Image.open(image_path)
    if image.mode != 'RGB':
        image = image.convert('RGB')
    image.save(output_path, 'PDF', resolution=100.0)
    return output_path

def convert_doc_to_pdf(doc_path):
    # Placeholder: You need libreoffice/unoconv on server for actual conversion
    output_path = os.path.join(OUTPUT_FOLDER, 'converted.pdf')
    os.system(f'libreoffice --headless --convert-to pdf "{doc_path}" --outdir "{OUTPUT_FOLDER}"')
    return output_path

def convert_ppt_to_pdf(ppt_path):
    # Placeholder: You need libreoffice/unoconv on server for actual conversion
    output_path = os.path.join(OUTPUT_FOLDER, 'converted.pdf')
    os.system(f'libreoffice --headless --convert-to pdf "{ppt_path}" --outdir "{OUTPUT_FOLDER}"')
    return output_path

def convert_excel_to_pdf(excel_path):
    # Placeholder: You need libreoffice/unoconv on server for actual conversion
    output_path = os.path.join(OUTPUT_FOLDER, 'converted.pdf')
    os.system(f'libreoffice --headless --convert-to pdf "{excel_path}" --outdir "{OUTPUT_FOLDER}"')
    return output_path

if __name__ == '__main__':
    app.run(debug=True, port=5000)
